import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches

# Create a new PowerPoint presentation
presentation = Presentation()

# Slide 1: Title Slide
slide = presentation.slides.add_slide(presentation.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "3-Month and 3-6 Month Plan for Data Engineering Team"
subtitle.text = "Presented by [Your Name] | [Date]"

# Slide 2: Agenda
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
title.text = "Agenda"

content = slide.placeholders[1]
content.text = (
    "1. Current Team and Challenges\n"
    "2. Proposed 3-Month Plan\n"
    "3. Proposed 3-6 Month Plan\n"
    "4. Expected Outcomes\n"
)

# Slide 3: Current Team and Challenges
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
title.text = "Current Team and Challenges"

content = slide.placeholders[1]
content.text = (
    "Team Overview:\n"
    "- Senior Manager: Analytics background\n"
    "- Senior Data Engineer: Strong in SQL and management\n"
    "- Junior Data Engineer: Good source system knowledge, weak coding\n"
    "- Offshore Team: Familiar with codebase, slow and inconsistent\n"
    "- New Member: Expertise in coding and architecture\n\n"
    "Key Challenges:\n"
    "- Poor code quality and technical debt\n"
    "- No intake process, reactive work handling\n"
    "- Skill gaps and inefficiencies\n"
)

# Slide 4: Proposed 3-Month Plan
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
title.text = "Proposed 3-Month Plan"

content = slide.placeholders[1]
content.text = (
    "1. Establish Intake Process:\n"
    "- Implement a request tracking system (e.g., Jira or Forms)\n"
    "- Weekly triage to prioritize tasks\n\n"
    "2. Address Technical Debt:\n"
    "- Refactor critical pipelines for maintainability\n"
    "- Document key parts of the codebase\n\n"
    "3. Upskill Team Members:\n"
    "- Training on coding and logic for Junior Data Engineer\n"
    "- Offshore team workshops on coding standards\n\n"
    "4. Implement Code Reviews:\n"
    "- Enforce quality checks for all code changes\n"
)

# Slide 5: Proposed 3-6 Month Plan
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
title.text = "Proposed 3-6 Month Plan"

content = slide.placeholders[1]
content.text = (
    "1. Build Modular Architecture:\n"
    "- Introduce reusable components for pipelines\n"
    "- Implement orchestration tools (e.g., Airflow)\n\n"
    "2. Automate Processes:\n"
    "- CI/CD pipelines for testing and deployment\n"
    "- Monitoring for data quality and pipeline performance\n\n"
    "3. Address Offshore Performance:\n"
    "- Evaluate and reassign responsibilities\n"
    "- Plan for hiring or restructuring, if needed\n\n"
    "4. Stakeholder Engagement:\n"
    "- Regular progress updates\n"
    "- Proactively manage expectations\n"
)

# Slide 6: Expected Outcomes
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
title.text = "Expected Outcomes"

content = slide.placeholders[1]
content.text = (
    "- Improved code quality and maintainability\n"
    "- Structured workflow with defined processes\n"
    "- Upskilled team members with reduced dependencies\n"
    "- Scalable architecture supporting future growth\n"
    "- Increased efficiency and stakeholder satisfaction\n"
)

# Save the presentation
file_path = "/mnt/data/Data_Engineering_Team_Plan.pptx"
presentation.save(file_path)

file_path
